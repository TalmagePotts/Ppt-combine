#!/usr/bin/env python3
"""
PowerPoint Combiner - GUI Version
Combine multiple PowerPoint and PDF files into one while preserving individual themes.
"""

import os
import sys
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import threading
import io
import copy
import subprocess

# Try to import pdf2image for PDF support
try:
    from pdf2image import convert_from_path
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False


class PowerPointCombinerGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Combiner")
        self.root.geometry("600x550")
        self.root.resizable(True, True)

        # Determine poppler path
        self.poppler_path = self.get_poppler_path()
        
        # Determine if we can use AppleScript (macOS + PPT installed)
        self.can_use_applescript = self.check_powerpoint_installed()

        # Variables
        self.input_folder = tk.StringVar()
        self.output_folder = tk.StringVar()
        self.output_filename = tk.StringVar(value="combined_presentation.pptx")
        self.convert_to_images = tk.BooleanVar(value=True)
        self.cancel_flag = threading.Event()

        self.setup_ui()
        
        if not PDF_SUPPORT:
            self.root.after(100, lambda: messagebox.showwarning(
                "PDF Support Missing", 
                "The 'pdf2image' library was not found.\nPDF files will be ignored."))
    
    def check_powerpoint_installed(self):
        """Check if Microsoft PowerPoint is available via AppleScript."""
        if sys.platform != 'darwin':
            return False
        try:
            # Check for PowerPoint bundle ID
            cmd = ['osascript', '-e', 'id of application "Microsoft PowerPoint"']
            result = subprocess.run(cmd, capture_output=True, text=True)
            return result.returncode == 0
        except:
            return False

    def get_poppler_path(self):
        """Get the path to bundled poppler binaries if frozen."""
        if getattr(sys, 'frozen', False):
            # If running as a PyInstaller bundle
            if hasattr(sys, '_MEIPASS'):
                # PyInstaller onefile or onedir
                bundled_path = os.path.join(sys._MEIPASS, 'poppler')
                if os.path.exists(bundled_path):
                    return bundled_path
            
            # Fallback for onedir if not in MEIPASS
            exe_dir = os.path.dirname(sys.executable)
            bundled_path = os.path.join(exe_dir, 'poppler')
            if os.path.exists(bundled_path):
                return bundled_path
                
        return None  # Rely on PATH

    def setup_ui(self):
        """Create the user interface."""
        # Main frame with padding
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))

        # Configure grid weights
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)

        # Title
        title_label = ttk.Label(main_frame, text="PowerPoint Combiner",
                                font=("Arial", 16, "bold"))
        title_label.grid(row=0, column=0, columnspan=3, pady=(0, 20))

        # Input folder selection
        ttk.Label(main_frame, text="Input Folder:").grid(row=1, column=0, sticky=tk.W, pady=5)
        ttk.Entry(main_frame, textvariable=self.input_folder, width=50).grid(
            row=1, column=1, sticky=(tk.W, tk.E), pady=5, padx=5)
        ttk.Button(main_frame, text="Browse...", command=self.browse_input_folder).grid(
            row=1, column=2, pady=5)

        # Output folder selection
        ttk.Label(main_frame, text="Output Folder:").grid(row=2, column=0, sticky=tk.W, pady=5)
        ttk.Entry(main_frame, textvariable=self.output_folder, width=50).grid(
            row=2, column=1, sticky=(tk.W, tk.E), pady=5, padx=5)
        ttk.Button(main_frame, text="Browse...", command=self.browse_output_folder).grid(
            row=2, column=2, pady=5)

        # Output filename
        ttk.Label(main_frame, text="Output Filename:").grid(row=3, column=0, sticky=tk.W, pady=5)
        ttk.Entry(main_frame, textvariable=self.output_filename, width=50).grid(
            row=3, column=1, sticky=(tk.W, tk.E), pady=5, padx=5)
            
        # Checkbox for Image Conversion
        ttk.Checkbutton(main_frame, text="Convert slides to Images (Preserves exact theme/fonts, not editable)", 
                       variable=self.convert_to_images).grid(row=4, column=0, columnspan=3, sticky=tk.W, pady=10)

        # Progress bar
        self.progress = ttk.Progressbar(main_frame, mode='indeterminate')
        self.progress.grid(row=5, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=20)

        # Status text area
        self.status_label = ttk.Label(main_frame, text="Status:")
        self.status_label.grid(row=6, column=0, sticky=tk.NW, pady=5)
        self.status_text = tk.Text(main_frame, height=10, width=60, state='disabled')
        self.status_text.grid(row=6, column=1, columnspan=2, sticky=(tk.W, tk.E), pady=5)

        # Scrollbar for status text
        scrollbar = ttk.Scrollbar(main_frame, orient=tk.VERTICAL, command=self.status_text.yview)
        scrollbar.grid(row=6, column=2, sticky=(tk.N, tk.S, tk.E), pady=5)
        self.status_text['yscrollcommand'] = scrollbar.set

        # Buttons Frame
        button_frame = ttk.Frame(main_frame)
        button_frame.grid(row=7, column=0, columnspan=3, pady=20)

        # Combine button
        self.combine_button = ttk.Button(button_frame, text="Combine Files",
                                         command=self.combine_powerpoints,
                                         style="Accent.TButton")
        self.combine_button.pack(side=tk.LEFT, padx=5)
        
        # Cancel button
        self.cancel_button = ttk.Button(button_frame, text="Cancel",
                                        command=self.cancel_operation,
                                        state='disabled')
        self.cancel_button.pack(side=tk.LEFT, padx=5)

        # Make rows expandable
        main_frame.rowconfigure(6, weight=1)

    def browse_input_folder(self):
        """Open dialog to select input folder."""
        folder = filedialog.askdirectory(title="Select folder containing PowerPoint/PDF files")
        if folder:
            self.input_folder.set(folder)
            # Auto-set output folder to same location if not set
            if not self.output_folder.get():
                self.output_folder.set(folder)

    def browse_output_folder(self):
        """Open dialog to select output folder."""
        folder = filedialog.askdirectory(title="Select output folder")
        if folder:
            self.output_folder.set(folder)

    def log_status(self, message):
        """Add message to status text area."""
        self.status_text.config(state='normal')
        self.status_text.insert(tk.END, message + "\n")
        self.status_text.see(tk.END)
        self.status_text.config(state='disabled')
        self.root.update()
        
    def cancel_operation(self):
        """Signal the operation to stop."""
        self.log_status("\nCancelling operation... please wait for current step to finish.")
        self.cancel_flag.set()
        self.cancel_button.config(state='disabled')

    def add_fitted_image_slide(self, prs, image_source, match_aspect_ratio=False):
        """Add a slide with an image. Optionally adjust presentation aspect ratio to match image."""
        layout_idx = 6 if len(prs.slide_layouts) > 6 else len(prs.slide_layouts) - 1
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        img_w = 0
        img_h = 0
        
        # Get dimensions
        if isinstance(image_source, (str, Path)):
             from PIL import Image
             with Image.open(str(image_source)) as img:
                 img_w, img_h = img.size
        else:
             # PIL Image object
             img_w, img_h = image_source.size
             
        if img_h == 0: return # Safety

        # Calculate Image Aspect Ratio
        img_ratio = img_w / img_h

        # Optionally resize presentation to match image aspect ratio
        # We keep width constant and adjust height
        if match_aspect_ratio:
            try:
                # Standard width is usually 10 inches or 13.33 inches (screen)
                # We trust the current width and adjust height
                new_slide_height = int(prs.slide_width / img_ratio)
                prs.slide_height = new_slide_height
                self.log_status(f"    Adjusted presentation aspect ratio to match image ({img_ratio:.2f})")
            except Exception as e:
                self.log_status(f"    Warning: Could not resize slide: {e}")

        # Calculate fit
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        slide_ratio = slide_width / slide_height
        
        if img_ratio > slide_ratio:
            # Image is wider -> Fit to width
            new_w = slide_width
            new_h = new_w / img_ratio
            left = 0
            top = (slide_height - new_h) / 2
        else:
            # Image is taller/same -> Fit to height
            new_h = slide_height
            new_w = new_h * img_ratio
            top = 0
            left = (slide_width - new_w) / 2
            
        # Add picture
        if isinstance(image_source, (str, Path)):
             slide.shapes.add_picture(str(image_source), left, top, new_w, new_h)
        else:
             # Stream from PIL image
             image_stream = io.BytesIO()
             image_source.save(image_stream, format='PNG')
             image_stream.seek(0)
             slide.shapes.add_picture(image_stream, left, top, new_w, new_h)

    def process_pdf(self, pdf_path, prs, is_first_file=False):
        """Convert PDF pages to images and add as slides."""
        if not PDF_SUPPORT:
            self.log_status(f"  Skipping {pdf_path.name}: pdf2image not installed.")
            return False

        self.log_status(f"  Converting PDF: {pdf_path.name}...")
        try:
            # Convert PDF to images
            try:
                images = convert_from_path(str(pdf_path), poppler_path=self.poppler_path)
            except Exception as e:
                if "poppler" in str(e).lower() or "pdfinfo" in str(e).lower():
                    self.log_status(f"  Error: Poppler is not installed/found.")
                    return False
                raise e

            for i, img in enumerate(images):
                # Resize slide to match the VERY FIRST image of the FIRST file
                should_match = (is_first_file and i == 0)
                self.add_fitted_image_slide(prs, img, match_aspect_ratio=should_match)
            
            self.log_status(f"  Added: {pdf_path.name} ({len(images)} slides)")
            return True

        except Exception as e:
            self.log_status(f"  Error processing PDF {pdf_path.name}: {str(e)}")
            return False

    def copy_slide_elements(self, source_slide, target_slide):
        """
        Copy shapes from source_slide to target_slide safely.
        """
        for shape in source_slide.shapes:
            try:
                # 1. Pictures
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if hasattr(shape, 'image'):
                        blob = shape.image.blob
                        image_stream = io.BytesIO(blob)
                        target_slide.shapes.add_picture(
                            image_stream, 
                            shape.left, shape.top, 
                            shape.width, shape.height
                        )
                
                # 2. Text Boxes and Shapes
                elif shape.has_text_frame:
                    new_shape = None
                    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        new_shape = target_slide.shapes.add_textbox(
                            shape.left, shape.top,
                            shape.width, shape.height
                        )
                    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        new_shape = target_slide.shapes.add_shape(
                            shape.auto_shape_type,
                            shape.left, shape.top,
                            shape.width, shape.height
                        )
                    
                    if new_shape and shape.text_frame.text:
                        new_shape.text_frame.text = shape.text_frame.text
                        # Try to copy simple paragraph formatting (bold, size, color)
                        # Note: This is partial support
                        try:
                            if shape.text_frame.paragraphs and new_shape.text_frame.paragraphs:
                                p_src = shape.text_frame.paragraphs[0]
                                p_dst = new_shape.text_frame.paragraphs[0]
                                if p_src.runs and p_dst.runs:
                                    r_src = p_src.runs[0]
                                    r_dst = p_dst.runs[0]
                                    if hasattr(r_src.font, 'size') and r_src.font.size:
                                        r_dst.font.size = r_src.font.size
                                    if hasattr(r_src.font, 'bold'):
                                        r_dst.font.bold = r_src.font.bold
                        except:
                            pass
                        
            except Exception as e:
                # Log to console if needed, but don't stop process
                print(f"Warning copying shape: {e}")
                pass

    def robust_xml_copy(self, source_slide, target_slide):
        """
        Attempt to copy XML but fix relationships to avoid corruption.
        This is experimental but better than raw copy.
        """
        # Copy all shapes from the original slide
        for shape in source_slide.shapes:
            try:
                # Get the shape element XML
                el = shape.element
                new_el = copy.deepcopy(el)
                
                # Check for relationships (images, etc)
                # This is a simplification; handling all rels correctly is complex
                # If we detect a relationship we can't handle easily, we skip the shape
                # to prevent corruption.
                
                # For now, we actually fallback to 'copy_slide_elements' for this shape
                # if it looks risky, or if we are in "Simple Mode".
                # But to truly support "Robust Copy", we would need to map rIds.
                
                # Given the complexity, we will stick to 'copy_slide_elements' (Reconstruction)
                # as the default 'Safe' method.
                pass
            except Exception:
                pass
                
    def convert_pptx_to_images_macos(self, input_path, output_folder):
        """Convert PPTX to Images by first exporting to PDF (Robust)."""
        import uuid
        
        input_abs = os.path.abspath(str(input_path))
        output_abs = os.path.abspath(str(output_folder))
        input_dir = os.path.dirname(input_abs)
        
        # Use a temp file in the SAME directory as input to ensure permissions match
        temp_filename = f"temp_export_{uuid.uuid4().hex[:8]}.pdf"
        temp_pdf = os.path.join(input_dir, temp_filename)
        
        # Ensure output directory exists
        os.makedirs(output_abs, exist_ok=True)
        
        # Create a placeholder file to ensure the path is valid for AppleScript
        try:
            with open(temp_pdf, 'wb') as f:
                pass
        except Exception as e:
            self.log_status(f"    Error creating placeholder file: {e}")
            return False
        
        # AppleScript to convert PPTX to PDF
        script = f'''
        tell application "Microsoft PowerPoint"
            -- activate
            try
                -- Open the file
                open POSIX file "{input_abs}" with read only
                set activePres to active presentation
                
                -- Save as PDF
                -- We use the existing placeholder file which helps with permissions/path resolution
                save activePres in POSIX file "{temp_pdf}" as save as PDF
                
                -- Close without saving changes
                close activePres saving no
                
                return "Success"
            on error errMsg
                try
                    close activePres saving no
                end try
                return "Error: " & errMsg
            end try
        end tell
        '''
        
        try:
            # 1. Export to PDF
            process = subprocess.run(
                ['/usr/bin/osascript', '-e', script], 
                capture_output=True, 
                text=True
            )
            
            if process.returncode != 0:
                self.log_status(f"    AppleScript System Error: {process.stderr.strip()}")
                return False
                
            result = process.stdout.strip()
            if result.startswith("Error:"):
                 self.log_status(f"    PPT Error: {result}")
                 # Cleanup placeholder if it failed
                 try: os.remove(temp_pdf) 
                 except: pass
                 return False
            
            # Check if file has size (meaning it was written to)
            if not os.path.exists(temp_pdf) or os.path.getsize(temp_pdf) == 0:
                self.log_status(f"    Error: PDF export failed (File empty or missing at {temp_pdf}).")
                return False
            
            # 2. Convert PDF to Images using pdf2image
            try:
                images = convert_from_path(temp_pdf, poppler_path=self.poppler_path)
                
                # Save images to output folder
                for i, img in enumerate(images):
                    img_path = os.path.join(output_abs, f"Slide_{i+1:03d}.png")
                    img.save(img_path, format='PNG')
                
                # Cleanup PDF
                try:
                    os.remove(temp_pdf)
                except:
                    pass
                return True
                
            except Exception as e:
                self.log_status(f"    Error converting intermediate PDF to images: {e}")
                return False

        except Exception as e:
            self.log_status(f"    Exception calling AppleScript: {str(e)}")
            return False

    def combine_powerpoints(self):
        """Combine files in a separate thread."""
        if not self.input_folder.get():
            messagebox.showerror("Error", "Please select an input folder")
            return

        if not self.output_folder.get():
            messagebox.showerror("Error", "Please select an output folder")
            return

        if not self.output_filename.get():
            messagebox.showerror("Error", "Please enter an output filename")
            return

        filename = self.output_filename.get()
        if not filename.endswith('.pptx'):
            filename += '.pptx'
            self.output_filename.set(filename)

        self.status_text.config(state='normal')
        self.status_text.delete(1.0, tk.END)
        self.status_text.config(state='disabled')
        
        # Reset status label
        self.status_label.config(text="Status:")

        # Reset cancel flag
        self.cancel_flag.clear()
        
        # Toggle buttons
        self.combine_button.config(state='disabled')
        self.cancel_button.config(state='normal')
        
        self.progress.start()

        thread = threading.Thread(target=self.do_combine, daemon=True)
        thread.start()

    def do_combine(self):
        """Perform the combination logic."""
        try:
            input_path = Path(self.input_folder.get())
            output_path = Path(self.output_folder.get()) / self.output_filename.get()
            convert_images = self.convert_to_images.get()

            # Get files
            files = []
            for ext in ["*.pptx", "*.pdf"]:
                files.extend(input_path.glob(ext))
            
            files = sorted(files, key=lambda p: p.name)
            files = [f for f in files if not f.name.startswith("~$")]

            if not files:
                self.root.after(0, lambda: messagebox.showerror("Error",
                    f"No PowerPoint or PDF files found in '{input_path}'"))
                return

            self.log_status(f"Found {len(files)} file(s):")
            for f in files:
                self.log_status(f"  - {f.name}")

            self.log_status("\nCombining presentations...")
            if convert_images:
                self.log_status("Mode: Converting slides to Images (Theme Preserved)")
                if not self.can_use_applescript:
                    self.log_status("Warning: Microsoft PowerPoint not found. PPTX conversion may fail or use text-only fallback.")
            else:
                self.log_status("Mode: Safe Copy (Theme may be lost, Text Editable)")

            combined_prs = None
            first_file = files[0]
            started_blank = False

            # Initialize base presentation
            if first_file.suffix.lower() == '.pptx':
                try:
                    combined_prs = Presentation(first_file)
                    self.log_status(f"  Added: {first_file.name} ({len(combined_prs.slides)} slides) [Base]")
                    start_index = 1
                    started_blank = False
                except Exception as e:
                    self.log_status(f"Error loading base file: {str(e)}")
                    raise e
            else:
                combined_prs = Presentation()
                start_index = 0
                started_blank = True

            # Process files
            files_to_process = files[start_index:]
            total_files = len(files_to_process)
            
            # Configure determinate progress bar
            self.root.after(0, lambda: self.progress.config(mode='determinate', maximum=total_files, value=0))

            for i, file_path in enumerate(files_to_process):
                # Check for cancellation
                if self.cancel_flag.is_set():
                    self.log_status("\nOperation cancelled by user.")
                    break

                current_num = i + 1
                
                # Update status label text
                self.root.after(0, lambda c=current_num, t=total_files: 
                               self.status_label.config(text=f"Status ({c}/{t}):"))
                
                self.log_status(f"\n[{current_num}/{total_files}] Processing: {file_path.name}...")
                
                is_very_first = (started_blank and i == 0)
                
                if file_path.suffix.lower() == '.pptx':
                    if convert_images and self.can_use_applescript:
                        # Convert PPTX to Images using AppleScript
                        self.log_status(f"  Converting PPTX to Images...")
                        
                        # Create a specific temp dir for this file
                        import shutil
                        temp_dir = input_path / f"temp_{file_path.stem}"
                        if temp_dir.exists():
                            shutil.rmtree(temp_dir)
                        os.makedirs(temp_dir)
                        
                        success = self.convert_pptx_to_images_macos(file_path, temp_dir)
                        
                        if success:
                            # Find all PNGs
                            image_files = list(temp_dir.glob("**/*.png")) + list(temp_dir.glob("**/*.PNG"))
                            
                            # Sort by slide number (Slide1, Slide2, ... Slide10)
                            # Natural sort
                            import re
                            def natural_sort_key(s):
                                return [int(c) if c.isdigit() else c.lower() for c in re.split(r'(\d+)', s.name)]
                            
                            image_files.sort(key=natural_sort_key)
                            
                            if not image_files:
                                self.log_status(f"    Warning: No images found after export.")
                                success = False
                            else:
                                # Add each image as a slide with proper fitting
                                for j, img_path in enumerate(image_files):
                                    if self.cancel_flag.is_set(): break # Check inner loop
                                    try:
                                        match = (is_very_first and j == 0)
                                        self.add_fitted_image_slide(combined_prs, img_path, match_aspect_ratio=match)
                                    except Exception as e:
                                        self.log_status(f"    Error adding image {img_path.name}: {e}")

                                self.log_status(f"  Added: {file_path.name} ({len(image_files)} slides)")
                                
                                # Cleanup
                                try:
                                    shutil.rmtree(temp_dir)
                                except:
                                    pass
                                
                                # Update progress
                                self.root.after(0, lambda val=current_num: self.progress.configure(value=val))
                                continue

                    # Fallback or Standard Mode
                    if convert_images:
                        self.log_status("    Warning: Falling back to Safe Copy.")

                    try:
                        prs = Presentation(file_path)
                        for slide in prs.slides:
                            if self.cancel_flag.is_set(): break
                            # Use blank layout from DESTINATION
                            layout_idx = 6 if len(combined_prs.slide_layouts) > 6 else len(combined_prs.slide_layouts) - 1
                            blank_layout = combined_prs.slide_layouts[layout_idx]
                            
                            new_slide = combined_prs.slides.add_slide(blank_layout)
                            
                            # Safe copy elements
                            self.copy_slide_elements(slide, new_slide)
                            
                        self.log_status(f"  Added: {file_path.name} ({len(prs.slides)} slides) [Text Only]")
                    except Exception as e:
                        self.log_status(f"  Error adding {file_path.name}: {str(e)}")

                elif file_path.suffix.lower() == '.pdf':
                    self.process_pdf(file_path, combined_prs, is_first_file=is_very_first)
                
                # Update progress
                self.root.after(0, lambda val=current_num: self.progress.configure(value=val))

            if not self.cancel_flag.is_set():
                combined_prs.save(str(output_path))

                self.log_status(f"\nSuccess! Created: {output_path}")
                self.log_status(f"Total slides: {len(combined_prs.slides)}")

                self.root.after(0, lambda: messagebox.showinfo("Success",
                    f"Successfully combined {len(files)} files!\n"
                    f"Output: {output_path}\n"
                    f"Total slides: {len(combined_prs.slides)}"))
            else:
                self.log_status(f"\nCancelled. Output not saved.")

        except Exception as e:
            self.log_status(f"\nError: {str(e)}")
            self.root.after(0, lambda: messagebox.showerror("Error",
                f"An error occurred:\n{str(e)}"))

        finally:
            self.root.after(0, lambda: self.combine_button.config(state='normal'))
            self.root.after(0, lambda: self.cancel_button.config(state='disabled'))
            self.root.after(0, lambda: self.progress.stop())


def main():
    """Main function to launch the GUI."""
    root = tk.Tk()
    app = PowerPointCombinerGUI(root)
    root.mainloop()


if __name__ == "__main__":
    main()