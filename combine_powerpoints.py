#!/usr/bin/env python3
"""
Combine multiple PowerPoint and PDF files into one PowerPoint presentation.
Source PowerPoint slides maintain their formatting.
PDF pages are converted to images and added as slides.
"""

import os
import sys
import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Try to import pdf2image for PDF support
try:
    from pdf2image import convert_from_path
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False


def process_pdf(pdf_path, prs):
    """
    Convert PDF pages to images and add them as slides to the presentation.
    
    Args:
        pdf_path: Path to the PDF file
        prs: The PowerPoint presentation object to add slides to
        
    Returns:
        bool: True if successful, False otherwise
    """
    if not PDF_SUPPORT:
        print(f"Error: pdf2image module not found. Cannot process {pdf_path.name}.")
        return False

    print(f"  Converting PDF: {pdf_path.name}...")
    try:
        # Convert PDF to images
        try:
            # images will be a list of PIL images
            images = convert_from_path(str(pdf_path))
        except Exception as e:
            if "poppler" in str(e).lower() or "pdfinfo" in str(e).lower():
                print(f"  Error: Poppler is not installed or not in PATH.")
                print(f"  Please install poppler to process PDFs.")
                print(f"  (e.g., 'brew install poppler' on macOS, 'apt-get install poppler-utils' on Linux)")
                return False
            raise e

        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for i, img in enumerate(images):
            # Use a blank layout (usually index 6)
            # If not available, use the last one (usually blank or simple)
            layout_idx = 6 if len(prs.slide_layouts) > 6 else len(prs.slide_layouts) - 1
            slide_layout = prs.slide_layouts[layout_idx]
            
            slide = prs.slides.add_slide(slide_layout)

            # Save image to bytes
            image_stream = io.BytesIO()
            img.save(image_stream, format='PNG')
            image_stream.seek(0)

            # Calculate fitting dimensions maintaining aspect ratio
            img_w, img_h = img.size
            aspect_ratio = img_w / img_h
            slide_ratio = slide_width / slide_height
            
            if aspect_ratio > slide_ratio:
                # Image is wider than slide (relative to aspect) -> Fit to width
                new_w = slide_width
                new_h = new_w / aspect_ratio
                left = 0
                top = (slide_height - new_h) / 2
            else:
                # Image is taller or same -> Fit to height
                new_h = slide_height
                new_w = new_h * aspect_ratio
                top = 0
                left = (slide_width - new_w) / 2
            
            slide.shapes.add_picture(image_stream, left, top, new_w, new_h)

        print(f"  Added: {pdf_path.name} ({len(images)} slides)")
        return True

    except Exception as e:
        print(f"  Error processing PDF {pdf_path.name}: {e}")
        return False


def copy_slide_elements(source_slide, target_slide):
    """
    Copy shapes from source_slide to target_slide safely.
    This avoids XML injection which can cause corruption.
    """
    for shape in source_slide.shapes:
        try:
            # 1. Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Extract image blob
                if hasattr(shape, 'image'):
                    blob = shape.image.blob
                    image_stream = io.BytesIO(blob)
                    target_slide.shapes.add_picture(
                        image_stream, 
                        shape.left, shape.top, 
                        shape.width, shape.height
                    )
            
            # 2. Text Boxes and Shapes
            elif shape.has_text_frame:
                new_shape = None
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    new_shape = target_slide.shapes.add_textbox(
                        shape.left, shape.top,
                        shape.width, shape.height
                    )
                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Try to preserve auto shape type
                    new_shape = target_slide.shapes.add_shape(
                        shape.auto_shape_type,
                        shape.left, shape.top,
                        shape.width, shape.height
                    )
                
                # Copy Text
                if new_shape and shape.text_frame.text:
                    new_shape.text_frame.text = shape.text_frame.text
                    # We could copy paragraphs/runs here for better fidelity
                    # but simple text is safest to start with.
                    
        except Exception as e:
            # If a specific shape fails, skip it but continue with others
            print(f"    Warning: Could not copy a shape: {e}")
            pass

def combine_powerpoints(input_folder, output_file):
    """
    Combine all PowerPoint and PDF files in a folder into a single presentation.

    Args:
        input_folder: Path to folder containing .pptx and .pdf files
        output_file: Path for the output combined .pptx file
    """
    input_path = Path(input_folder)

    if not input_path.exists():
        print(f"Error: Folder '{input_folder}' does not exist.")
        return False

    # Get all PowerPoint and PDF files
    files = []
    # Glob both extensions
    for ext in ["*.pptx", "*.pdf"]:
        files.extend(input_path.glob(ext))
    
    # Sort files by name
    files = sorted(files, key=lambda p: p.name)

    # Filter out temporary files
    files = [f for f in files if not f.name.startswith("~$")]

    if not files:
        print(f"No PowerPoint or PDF files found in '{input_folder}'")
        return False

    print(f"Found {len(files)} file(s):")
    for f in files:
        print(f"  - {f.name}")

    print(f"\nCombining presentations...")
    
    combined_prs = None
    
    # Determine base presentation
    first_file = files[0]
    
    if first_file.suffix.lower() == '.pptx':
        try:
            combined_prs = Presentation(first_file)
            print(f"  Added: {first_file.name} ({len(combined_prs.slides)} slides) [Base]")
            # We have already added the first file
            start_index = 1
        except Exception as e:
            print(f"Error loading base file {first_file.name}: {e}")
            return False
    else:
        # First file is PDF (or other if we supported it), start with blank
        combined_prs = Presentation()
        # We need to process the first file since it wasn't a base PPTX
        start_index = 0

    # Process remaining files (or all if first was PDF)
    for file_path in files[start_index:]:
        if file_path.suffix.lower() == '.pptx':
            try:
                prs = Presentation(file_path)

                for slide in prs.slides:
                    # Use a blank layout from the DESTINATION presentation
                    # Index 6 is usually "Blank" in standard templates
                    layout_idx = 6 if len(combined_prs.slide_layouts) > 6 else len(combined_prs.slide_layouts) - 1
                    blank_layout = combined_prs.slide_layouts[layout_idx]
                    
                    new_slide = combined_prs.slides.add_slide(blank_layout)
                    
                    # Safe copy of elements
                    copy_slide_elements(slide, new_slide)

                print(f"  Added: {file_path.name} ({len(prs.slides)} slides)")
            except Exception as e:
                print(f"  Error adding {file_path.name}: {e}")
        
        elif file_path.suffix.lower() == '.pdf':
            process_pdf(file_path, combined_prs)

    # Save the combined presentation
    try:
        combined_prs.save(output_file)
        print(f"\nSuccessfully created: {output_file}")
        print(f"Total slides: {len(combined_prs.slides)}")
        return True
    except Exception as e:
        print(f"Error saving output file: {e}")
        return False


def main():
    """Main function to handle command line arguments."""
    if len(sys.argv) < 2:
        # Use current directory if no argument provided
        input_folder = "."
        output_file = "combined_presentation.pptx"
    elif len(sys.argv) == 2:
        input_folder = sys.argv[1]
        output_file = "combined_presentation.pptx"
    else:
        input_folder = sys.argv[1]
        output_file = sys.argv[2]

    print(f"Input folder: {input_folder}")
    print(f"Output file: {output_file}\n")

    success = combine_powerpoints(input_folder, output_file)

    if success:
        sys.exit(0)
    else:
        sys.exit(1)


if __name__ == "__main__":
    main()
